import os
from datetime import datetime

import openpyxl
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor

from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches as PptInches, Pt as PptPt

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS_DIR = os.path.join(PROJECT_ROOT, "docs")
SCREENSHOTS_DIR = os.path.join(PROJECT_ROOT, "screenshots")
OUTPUT_DIR = DOCS_DIR

AZUL_PRINCIPAL = RGBColor(0x12, 0x3E, 0x7A)
AZUL_ACCENT = RGBColor(0x0D, 0x6E, 0xFD)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRIS = RGBColor(0x66, 0x66, 0x66)
NEGRO = RGBColor(0x00, 0x00, 0x00)

AZUL_PPT = PptRGBColor(0x12, 0x3E, 0x7A)
BLANCO_PPT = PptRGBColor(0xFF, 0xFF, 0xFF)
GRIS_PPT = PptRGBColor(0xF4, 0xF8, 0xFC)
AMARILLO_PPT = PptRGBColor(0xFF, 0xC1, 0x07)
VERDE_PPT = PptRGBColor(0x2E, 0xCC, 0x71)
NEGRO_PPT = PptRGBColor(0x00, 0x00, 0x00)
GRIS_CLARO_PPT = PptRGBColor(0xBB, 0xBB, 0xBB)
GRIS_OSCURO_PPT = PptRGBColor(0x66, 0x66, 0x66)

SCREENSHOT_MAP = {
    "full_page_load": "full_page_load.png",
    "map_area": "map_area_detail.png",
    "chart": "chart_section.png",
    "tomtom_flow": "tomtom_flow_map.png",
    "tomtom_incidents": "tomtom_fullpage_incidents43.png",
    "tomtom_stats": "tomtom-stats.png",
    "login": "page-top-viewport.png",
    "dashboard": "full-page-initial.png",
}


def get_screenshot(name):
    filename = SCREENSHOT_MAP.get(name)
    if filename:
        path = os.path.join(SCREENSHOTS_DIR, filename)
        if os.path.exists(path):
            return path
    alt_path = os.path.join(SCREENSHOTS_DIR, name)
    if os.path.exists(alt_path):
        return alt_path
    return None


def add_screenshot_to_cell(cell, img_path, width=Inches(2.5)):
    if img_path and os.path.exists(img_path):
        try:
            cell.paragraphs[0].clear()
            run = cell.paragraphs[0].add_run()
            run.add_picture(img_path, width=width)
        except Exception:
            pass


# ──────────────────────────────────────────────
#  1. FICHA TÉCNICA (Word)
# ──────────────────────────────────────────────
def create_ficha_tecnica():
    doc = Document()

    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(11)
    font.color.rgb = NEGRO

    for level in range(1, 4):
        hs = doc.styles[f"Heading {level}"]
        hs.font.color.rgb = AZUL_PRINCIPAL
        hs.font.bold = True

    # Title
    title = doc.add_heading("Ficha Técnica del Proyecto", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.runs[0]
    run.font.color.rgb = AZUL_PRINCIPAL
    run.font.size = Pt(26)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run("Medellín Movilidata OS — HackData CTGI SENA 2026")
    run.font.size = Pt(14)
    run.font.color.rgb = AZUL_ACCENT

    doc.add_paragraph()

    # Info table
    info_data = [
        ("Nombre del Proyecto", "Medellín Movilidata OS"),
        ("Evento", "HackData CTGI SENA 2026 — Reto de Movilidad"),
        ("Ciudad", "Medellín, Colombia"),
        ("Versión", "1.0.0 (final)"),
        ("Fecha de elaboración", datetime.now().strftime("%d/%m/%Y")),
        ("Repositorio", "Git — rama final"),
        ("Frontend principal", "Vue 3 + Vite 8 + Bootstrap 5.3 + Leaflet + Chart.js + ECharts"),
        ("Frontend migración", "React 18 + Vite 5 + Tailwind CSS 4 + Leaflet + Chart.js"),
        ("Backend", "Django 5.1+ / Django REST Framework 3.15+"),
        ("Base de datos", "SQLite (desarrollo) / MySQL (producción)"),
        ("Motor de mapas", "Leaflet + leaflet.heat + TomTom Traffic API"),
        ("API externas", "TomTom Traffic, SIATA (Medellín clima), OpenWeatherMap"),
        ("Autenticación", "DRF Token Auth (rest_framework.authtoken)"),
        ("ML Predicción", "Regresión lineal con pesos temporales y factores día-semana"),
        ("PWA", "vite-plugin-pwa con service worker y offline.html"),
        ("Python", "3.11+"),
        ("Node.js", "20+, npm 10+"),
    ]

    table = doc.add_table(rows=len(info_data), cols=2)
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    for i, (key, val) in enumerate(info_data):
        row = table.rows[i]
        row.cells[0].text = key
        row.cells[1].text = val
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.style.font.size = Pt(10)
                for run in paragraph.runs:
                    run.font.size = Pt(10)

    doc.add_paragraph()

    # Architecture
    doc.add_heading("Arquitectura del Sistema", level=1)
    doc.add_paragraph(
        "El sistema sigue una arquitectura frontend-backend separada:"
    )
    doc.add_paragraph(
        "El frontend Vue 3 consume la API REST de Django a través de fetch(). "
        "El mapa Leaflet se renderiza en el componente Inicio.vue con capas de calor, "
        "polígonos de zonas de riesgo, tráfico TomTom e incidentes. "
        "El backend sirve datos desde SQLite con 3 modelos principales: "
        "Accident, Zone y WeatherRecord."
    )

    # Components
    doc.add_heading("Componentes del Frontend (Vue 3)", level=2)

    comps = [
        ("Inicio.vue", "Dashboard principal: mapa Leaflet, heatmap, zonas, clima, gráfica Chart.js, predicción ML"),
        ("AccidentReporter.vue", "Mapa con capa de calor, filtro por severidad, formulario click-to-report"),
        ("Login.vue", "Autenticación token, panel admin con CRUD de accidentes/usuarios"),
        ("RealtimeTracker.vue", "GPS tracking, rutas simuladas MIO/Metro, TomTom tráfico, capa MapGIS"),
        ("Nosotros.vue", "Misión, visión y equipo"),
        ("Servicios.vue", "Seis tarjetas de servicio: visualización, alertas, analítica, API, IoT, soporte"),
        ("Contacto.vue", "Formulario de contacto con validación"),
    ]
    comp_table = doc.add_table(rows=len(comps) + 1, cols=2)
    comp_table.style = "Light Grid Accent 1"
    header = comp_table.rows[0]
    header.cells[0].text = "Componente"
    header.cells[1].text = "Descripción"
    for i, (name, desc) in enumerate(comps, 1):
        comp_table.rows[i].cells[0].text = name
        comp_table.rows[i].cells[1].text = desc

    for row in comp_table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)

    doc.add_paragraph()

    # API Endpoints
    doc.add_heading("Endpoints de la API", level=2)

    endpoints = [
        ("GET /api/accidents/", "Lista todos los accidentes"),
        ("GET /api/accidents/?hour_from=&hour_to=", "Filtra accidentes por rango horario"),
        ("GET /api/zones/", "Lista zonas de riesgo con geometría GeoJSON"),
        ("GET /api/weather/", "Estado del clima simulado"),
        ("POST /api/simulate_rain/", "Alterna el estado de lluvia simulado"),
        ("GET /api/congestion_prediction/", "Predicción ML de congestión (próximas 2h)"),
        ("GET /api/congestion_prediction/?hour=N", "Predicción para una hora específica"),
    ]
    ep_table = doc.add_table(rows=len(endpoints) + 1, cols=2)
    ep_table.style = "Light Grid Accent 1"
    header = ep_table.rows[0]
    header.cells[0].text = "Endpoint"
    header.cells[1].text = "Descripción"
    for i, (ep, desc) in enumerate(endpoints, 1):
        ep_table.rows[i].cells[0].text = ep
        ep_table.rows[i].cells[1].text = desc

    for row in ep_table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)

    doc.add_paragraph()

    # Models
    doc.add_heading("Modelo de Datos", level=2)
    doc.add_paragraph("Accident (lat, lng, intensity, hour, date)")
    doc.add_paragraph("Zone (name, risk_level, geometry — GeoJSON)")
    doc.add_paragraph("WeatherRecord (location, condition, temperature, humidity, pressure, wind_speed, is_raining, source, recorded_at)")

    # Screenshot
    img = get_screenshot("full_page_load")
    if img:
        doc.add_paragraph()
        doc.add_heading("Captura de Pantalla", level=2)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(img, width=Inches(5.5))

    # Save
    path = os.path.join(OUTPUT_DIR, "Ficha_Tecnica_Medellin_Movilidata_OS.docx")
    doc.save(path)
    print(f"✅ Ficha técnica creada: {path}")
    return path


# ──────────────────────────────────────────────
#  2. MATRIZ DE REQUERIMIENTOS (Excel)
# ──────────────────────────────────────────────
def create_matriz_requerimientos():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Matriz de Requerimientos"

    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="123E7A", end_color="123E7A", fill_type="solid")
    header_alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    cell_alignment = Alignment(vertical="top", wrap_text=True)

    headers = ["ID", "Módulo", "Requerimiento", "Descripción", "Prioridad", "Estado", "Tipo"]
    col_widths = [8, 18, 30, 55, 12, 14, 14]

    for col_num, (header, width) in enumerate(zip(headers, col_widths), 1):
        cell = ws.cell(row=1, column=col_num, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_alignment
        ws.column_dimensions[get_column_letter(col_num)].width = width

    reqs = [
        # Dashboard
        ("REQ-001", "Dashboard", "Visualización de mapa interactivo",
         "El sistema debe mostrar un mapa Leaflet centrado en Medellín con capacidad de zoom y capas superpuestas.",
         "Alta", "Implementado", "Funcional"),
        ("REQ-002", "Dashboard", "Heatmap de accidentes",
         "El mapa debe renderizar un heatmap con los accidentes registrados, usando leaflet.heat.",
         "Alta", "Implementado", "Funcional"),
        ("REQ-003", "Dashboard", "Polígonos de zonas de riesgo",
         "El mapa debe mostrar polígonos coloreados por nivel de riesgo (alta/media/baja) desde la API /api/zones/.",
         "Alta", "Implementado", "Funcional"),
        ("REQ-004", "Dashboard", "Gráfica de tendencia horaria",
         "Dashboard debe incluir un gráfico de barras Chart.js con accidentes e intensidad acumulada por hora.",
         "Media", "Implementado", "Funcional"),
        ("REQ-005", "Dashboard", "Widget de clima",
         "Dashboard debe mostrar clima actual (temperatura, condición, lluvia) desde SIATA o API interna.",
         "Media", "Implementado", "Funcional"),
        ("REQ-006", "Dashboard", "Widget de predicción ML",
         "Dashboard debe predecir congestión para las próximas 2 horas usando regresión lineal con datos históricos.",
         "Media", "Implementado", "Funcional"),

        # Accident Reporter
        ("REQ-007", "Accidentes", "Registro de accidentes",
         "El usuario debe poder reportar un accidente haciendo clic en el mapa y llenando un formulario.",
         "Alta", "Implementado", "Funcional"),
        ("REQ-008", "Accidentes", "Filtro por severidad",
         "El mapa de accidentes debe poder filtrarse por nivel de severidad (baja/media/alta).",
         "Media", "Implementado", "Funcional"),

        # Tracking
        ("REQ-009", "Tracker", "Seguimiento GPS en tiempo real",
         "El módulo RealtimeTracker debe obtener la ubicación GPS del dispositivo usando Geolocation API.",
         "Alta", "Implementado", "Funcional"),
        ("REQ-010", "Tracker", "Visualización de tráfico TomTom",
         "El tracker debe mostrar segmentos de flujo vehicular e incidentes desde la API TomTom Traffic.",
         "Alta", "Implementado", "Funcional"),
        ("REQ-011", "Tracker", "Rutas simuladas",
         "El tracker debe simular rutas de MIO bus, Metro y vehículo particular sobre el mapa.",
         "Media", "Implementado", "Funcional"),

        # Auth / Admin
        ("REQ-012", "Autenticación", "Inicio de sesión con token",
         "El admin debe autenticarse con username/password y recibir un token DRF almacenado en localStorage.",
         "Alta", "Implementado", "Funcional"),
        ("REQ-013", "Autenticación", "CRUD de accidentes",
         "El panel admin debe permitir crear, leer, actualizar y eliminar accidentes.",
         "Alta", "Implementado", "Funcional"),
        ("REQ-014", "Autenticación", "CRUD de usuarios",
         "El panel admin debe permitir listar y eliminar usuarios del sistema.",
         "Media", "Implementado", "Funcional"),

        # Weather
        ("REQ-015", "Clima", "Simulación de lluvia",
         "El sistema debe permitir activar/desactivar lluvia simulada desde el frontend.",
         "Media", "Implementado", "Funcional"),
        ("REQ-016", "Clima", "Integración SIATA",
         "El frontend debe consultar datos meteorológicos reales desde la API SIATA de Medellín.",
         "Media", "Implementado", "Funcional"),

        # ML
        ("REQ-017", "ML", "Predicción de congestión",
         "El endpoint /api/congestion_prediction/ debe devolver predicción usando regresión lineal ponderada.",
         "Media", "Implementado", "Funcional"),

        # PWA
        ("REQ-018", "PWA", "Service worker offline",
         "La aplicación debe registrar un service worker que cachee recursos para funcionar offline.",
         "Baja", "Implementado", "Funcional"),
        ("REQ-019", "PWA", "Página offline",
         "Debe existir una página offline.html de respaldo cuando no hay conexión.",
         "Baja", "Implementado", "Funcional"),

        # Tech
        ("REQ-020", "Backend", "API REST con Django REST Framework",
         "El backend debe exponer una API REST con endpoints para accidentes, zonas, clima y predicción.",
         "Alta", "Implementado", "Técnico"),
        ("REQ-021", "Backend", "Carga de datos desde JSON",
         "El backend debe poder cargar datos semilla desde archivos JSON con el comando load_data.",
         "Alta", "Implementado", "Técnico"),
        ("REQ-022", "Backend", "Parseo de GeoJSON",
         "El sistema debe incluir un script para parsear archivos GeoJSON de incidentes de tránsito.",
         "Media", "Implementado", "Técnico"),
        ("REQ-023", "Backend", "CORS habilitado",
         "El backend debe permitir CORS para los orígenes localhost:5173 y localhost:8000.",
         "Alta", "Implementado", "Técnico"),
        ("REQ-024", "Frontend", "Mapbox GL como motor primario con fallback a Leaflet",
         "Inicio.vue debe intentar cargar Mapbox GL primero; si falla, usar Leaflet como alternativa.",
         "Alta", "Implementado", "Técnico"),
        ("REQ-025", "Documentación", "Manual de usuario",
         "Debe existir un manual de usuario que explique cómo usar el sistema paso a paso.",
         "Media", "Implementado", "Documentación"),
        ("REQ-026", "Documentación", "Ficha técnica",
         "Debe existir una ficha técnica con especificaciones del sistema.",
         "Media", "Implementado", "Documentación"),
    ]

    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    priority_colors = {
        "Alta": PatternFill(start_color="FFCCCC", end_color="FFCCCC", fill_type="solid"),
        "Media": PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid"),
        "Baja": PatternFill(start_color="D9EAD3", end_color="D9EAD3", fill_type="solid"),
    }
    state_colors = {
        "Implementado": PatternFill(start_color="D9EAD3", end_color="D9EAD3", fill_type="solid"),
        "Pendiente": PatternFill(start_color="F4CCCC", end_color="F4CCCC", fill_type="solid"),
    }

    for row_num, (req_id, module, name, desc, priority, state, rtype) in enumerate(reqs, 2):
        data = [req_id, module, name, desc, priority, state, rtype]
        for col_num, value in enumerate(data, 1):
            cell = ws.cell(row=row_num, column=col_num, value=value)
            cell.alignment = cell_alignment
            cell.border = thin_border
            cell.font = Font(name="Calibri", size=10)
        # Color priority column
        p_cell = ws.cell(row=row_num, column=5)
        if priority in priority_colors:
            p_cell.fill = priority_colors[priority]
        s_cell = ws.cell(row=row_num, column=6)
        if state in state_colors:
            s_cell.fill = state_colors[state]

    ws.auto_filter.ref = f"A1:G{len(reqs) + 1}"
    ws.freeze_panes = "A2"

    # Row height
    for row in range(2, len(reqs) + 2):
        ws.row_dimensions[row].height = 40

    path = os.path.join(OUTPUT_DIR, "Matriz_Requerimientos_Medellin_Movilidata_OS.xlsx")
    wb.save(path)
    print(f"✅ Matriz de requerimientos creada: {path}")
    return path


# ──────────────────────────────────────────────
#  3. PRESENTACIÓN (PowerPoint)
# ──────────────────────────────────────────────
def create_presentacion():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    def add_bg(slide, color=AZUL_PPT):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=BLANCO_PPT, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
        txBox = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptPt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def add_img_safe(slide, img_name, left, top, width=None, height=None):
        path = os.path.join(SCREENSHOTS_DIR, img_name)
        if not os.path.exists(path):
            path = os.path.join(SCREENSHOTS_DIR, SCREENSHOT_MAP.get(img_name, ""))
        if os.path.exists(path):
            kw = {}
            if width:
                kw["width"] = PptInches(width)
            if height:
                kw["height"] = PptInches(height)
            slide.shapes.add_picture(path, PptInches(left), PptInches(top), **kw)
            return True
        return False

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide1)
    add_textbox(slide1, 1.5, 1.5, 10, 1.5, "Medellín Movilidata OS", font_size=48, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, 1.5, 3.2, 10, 1, "Plataforma de Movilidad Segura Guiada por Datos", font_size=24, color=AMARILLO_PPT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, 1.5, 4.5, 10, 0.8, "HackData CTGI SENA 2026 — Reto de Movilidad", font_size=18, color=BLANCO_PPT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, 1.5, 5.5, 10, 0.6, f"Medellín, Colombia — {datetime.now().strftime('%B %Y')}", font_size=14, color=GRIS_CLARO_PPT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide1, 5.5, 6.5, 3, 0.5, "Equipo Medellín Data", font_size=14, color=BLANCO_PPT, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Problem ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide2, PptRGBColor(0xF4, 0xF8, 0xFC))
    add_textbox(slide2, 0.5, 0.3, 12, 1, "¿Cuál es el problema?", font_size=36, bold=True, color=AZUL_PPT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide2, 1, 1.8, 11, 5.5,
        "En Medellín, los accidentes de tránsito son una de las principales causas de mortalidad evitable.\n\n"
        "Los ciudadanos carecen de herramientas digitales que integren:\n\n"
        "  •  Datos históricos de accidentes\n"
        "  •  Visualización en tiempo real del tráfico\n"
        "  •  Predicción de zonas de alto riesgo\n"
        "  •  Alertas climáticas y de congestión\n\n"
        "Actualmente, la información está fragmentada entre múltiples fuentes (SIATA, TomTom, "
        "Secretaría de Movilidad) sin una interfaz unificada que permita tomar decisiones informadas.",
        font_size=18, color=NEGRO_PPT, alignment=PP_ALIGN.LEFT)

    # ── Slide 3: Solution ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide3)
    add_textbox(slide3, 0.5, 0.3, 12, 1, "Nuestra Solución", font_size=36, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide3, 1, 1.5, 11, 5.5,
        "Medellín Movilidata OS es una plataforma de código abierto que integra:\n\n"
        "  🗺️  Mapa interactivo con heatmap de accidentes y zonas de riesgo\n"
        "  🌧️  Datos climáticos en tiempo real (SIATA + OpenWeather)\n"
        "  🚦  Visualización de tráfico TomTom con incidentes y flujo vehicular\n"
        "  🤖  Predicción ML de congestión usando regresión lineal ponderada\n"
        "  📱  Aplicación web progresiva (PWA) con soporte offline\n"
        "  🔐  Panel administrativo con autenticación y CRUD\n"
        "  📊  Dashboard con gráficas interactivas (Chart.js + ECharts)",
        font_size=18, color=BLANCO_PPT, alignment=PP_ALIGN.LEFT)

    # ── Slide 4: Architecture ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide4, PptRGBColor(0xF4, 0xF8, 0xFC))
    add_textbox(slide4, 0.5, 0.3, 12, 1, "Arquitectura Técnica", font_size=36, bold=True, color=AZUL_PPT, alignment=PP_ALIGN.CENTER)

    arch_items = [
        ("Frontend Vue 3", "Vite 8 · Bootstrap 5\nLeaflet + Mapbox GL\nChart.js + ECharts\nTomTom Traffic API", 0.5, 1.5, 3.5, 3),
        ("Backend Django", "Django 5.1 + DRF 3.15\nSQLite / MySQL\nToken Auth\n5 endpoints REST", 4.8, 1.5, 3.5, 3),
        ("ML Engine", "scikit-learn\nRegresión Lineal\nPesos temporales\nFactores día-semana", 9.2, 1.5, 3.5, 3),
    ]
    for text, subtext, left, top, w, h in arch_items:
        shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(left), PptInches(top), PptInches(w), PptInches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = AZUL_PPT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = PptPt(20)
        run.font.bold = True
        run.font.color.rgb = BLANCO_PPT
        run.font.name = "Calibri"
        tf.add_paragraph().text = ""
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtext
        run2.font.size = PptPt(14)
        run2.font.color.rgb = AMARILLO_PPT
        run2.font.name = "Calibri"

    add_textbox(slide4, 0.5, 5.2, 12, 2,
        "APIs Externas: TomTom Traffic · SIATA (Medellín) · OpenWeatherMap\n"
        "Cliente: Navegador Web (Chrome/Firefox) · PWA con Service Worker · Geolocalización GPS",
        font_size=14, color=GRIS_OSCURO_PPT, alignment=PP_ALIGN.CENTER)

    # ── Slide 5: Screenshot Map ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide5, PptRGBColor(0xF4, 0xF8, 0xFC))
    add_textbox(slide5, 0.5, 0.3, 12, 0.8, "Dashboard — Mapa y Zonas de Riesgo", font_size=28, bold=True, color=AZUL_PPT, alignment=PP_ALIGN.CENTER)
    if not add_img_safe(slide5, "full-page-initial.png", 0.8, 1.3, width=11.5):
        add_img_safe(slide5, "full_page_load.png", 0.8, 1.3, width=11.5)

    # ── Slide 6: Screenshot Charts ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide6, PptRGBColor(0xF4, 0xF8, 0xFC))
    add_textbox(slide6, 0.5, 0.3, 12, 0.8, "Dashboard — Gráficas y Predicción ML", font_size=28, bold=True, color=AZUL_PPT, alignment=PP_ALIGN.CENTER)
    add_img_safe(slide6, "chart_section.png", 0.5, 1.3, width=6)
    add_img_safe(slide6, "tomtom-stats.png", 6.8, 1.3, width=6)

    # ── Slide 7: TomTom ──
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide7, PptRGBColor(0xF4, 0xF8, 0xFC))
    add_textbox(slide7, 0.5, 0.3, 12, 0.8, "Tráfico en Tiempo Real — TomTom Traffic", font_size=28, bold=True, color=AZUL_PPT, alignment=PP_ALIGN.CENTER)
    add_img_safe(slide7, "tomtom_flow_map.png", 0.5, 1.3, width=6)
    add_img_safe(slide7, "tomtom_fullpage_incidents43.png", 6.8, 1.3, width=5.8)

    # ── Slide 8: ML ──
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide8, PptRGBColor(0xF4, 0xF8, 0xFC))
    add_textbox(slide8, 0.5, 0.3, 12, 0.8, "Motor de Predicción ML", font_size=28, bold=True, color=AZUL_PPT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, 0.8, 1.5, 11.5, 5,
        "Algoritmo: Regresión Lineal con scikit-learn\n\n"
        "Características:\n"
        "  •  Pesos exponenciales: datos recientes (últimos 3 días) tienen 50% del peso total\n"
        "  •  Factor día-semana: fines de semana tienen 20% menos congestión\n"
        "  •  Detección de tendencia: pendiente de regresión lineal sobre datos históricos\n"
        "  •  Intervalos de confianza: entre 0.3 y 0.95 según R², muestras y horizonte\n"
        "  •  Niveles de riesgo: baja (<4 accidentes), media (4-7), alta (>=8)\n\n"
        "Endpoint: GET /api/congestion_prediction/?hour=N\n"
        "Respuesta: predicción para las próximas 2 horas con nivel de riesgo y confianza",
        font_size=16, color=NEGRO_PPT, alignment=PP_ALIGN.LEFT)

    # ── Slide 9: Tech Stack ──
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide9, PptRGBColor(0xF4, 0xF8, 0xFC))
    add_textbox(slide9, 0.5, 0.3, 12, 0.8, "Stack Tecnológico", font_size=28, bold=True, color=AZUL_PPT, alignment=PP_ALIGN.CENTER)

    stack_data = [
        ("Frontend", "Vue 3 · Vite 8 · Bootstrap 5\nLeaflet · Mapbox GL\nChart.js · ECharts\nPWA (vite-plugin-pwa)"),
        ("Backend", "Django 5.1+ · DRF 3.15+\nSQLite / MySQL\nToken Auth · CORS\nPython 3.11+"),
        ("ML & Datos", "scikit-learn · NumPy\nRegresión Lineal\nGeoJSON · Excel parsing\nSIATA + OpenWeather"),
        ("DevOps", "Git · GitHub\nVite · Hot Reload\nPWA Service Worker\nnpm / pip"),
    ]
    for i, (title, items) in enumerate(stack_data):
        left = 0.5 + i * 3.2
        shape = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(left), PptInches(1.5), PptInches(3), PptInches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = AZUL_PPT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = PptPt(20)
        run.font.bold = True
        run.font.color.rgb = BLANCO_PPT
        run.font.name = "Calibri"
        tf.add_paragraph().text = ""
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = items
        run2.font.size = PptPt(14)
        run2.font.color.rgb = AMARILLO_PPT
        run2.font.name = "Calibri"

    # ── Slide 10: Timeline ──
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide10)
    add_textbox(slide10, 0.5, 0.3, 12, 0.8, "Roadmap del Proyecto", font_size=36, bold=True, alignment=PP_ALIGN.CENTER)
    timeline = [
        ("Fase 1", "MVP Core", "Mapa Leaflet, heatmap, API REST, carga datos"),
        ("Fase 2", "Tracking", "GPS, TomTom Traffic, rutas simuladas"),
        ("Fase 3", "ML Prediction", "Regresión lineal, pesos, endpoint predicción"),
        ("Fase 4", "PWA & Admin", "Service worker, offline, CRUD, autenticación"),
        ("Fase 5", "Documentación", "Manuales, ficha técnica, matriz req, presentación"),
    ]
    for i, (phase, title, desc) in enumerate(timeline):
        y = 1.8 + i * 1.1
        shape = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(1), PptInches(y), PptInches(11), PptInches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = AMARILLO_PPT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = f"  {phase}: {title} — {desc}"
        run.font.size = PptPt(16)
        run.font.bold = True
        run.font.color.rgb = AZUL_PPT
        run.font.name = "Calibri"

    # ── Slide 11: Thank You ──
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide11)
    add_textbox(slide11, 1.5, 2, 10, 1.5, "¡Gracias!", font_size=56, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide11, 1.5, 3.8, 10, 1, "Medellín Movilidata OS", font_size=28, color=AMARILLO_PPT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide11, 1.5, 5, 10, 0.8, "Movilidad segura guiada por datos", font_size=20, color=BLANCO_PPT, alignment=PP_ALIGN.CENTER)

    path = os.path.join(OUTPUT_DIR, "Presentacion_Medellin_Movilidata_OS.pptx")
    prs.save(path)
    print(f"✅ Presentación creada: {path}")
    return path


# ──────────────────────────────────────────────
#  Main
# ──────────────────────────────────────────────
if __name__ == "__main__":
    print("Generando documentación del proyecto...\n")
    create_ficha_tecnica()
    print()
    create_matriz_requerimientos()
    print()
    create_presentacion()
    print("\n✅ Todos los documentos generados correctamente en docs/")
